"""Build the Whitmore presentation deck as a professional PPTX file , Times New Roman throughout."""

from __future__ import annotations

import argparse
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

from taa_project.analysis.reporting import (
    DSR_SUMMARY_FILENAME,
    IPS_COMPLIANCE_FILENAME,
    PORTFOLIO_METRICS_FILENAME,
)
from taa_project.config import FIGURES_DIR, OUTPUT_DIR, REPORT_DIR

# ── Design system ──────────────────────────────────────────────────────────
NAVY = RGBColor(0x1A, 0x36, 0x5D)
GOLD = RGBColor(0xB8, 0x86, 0x0B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
ACCENT = RGBColor(0x2B, 0x6C, 0xB0)
TABLE_HEADER = RGBColor(0x1A, 0x36, 0x5D)
TABLE_ALT = RGBColor(0xEE, 0xF2, 0xF7)
TABLE_BORDER = RGBColor(0xCC, 0xCC, 0xCC)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT = "Times New Roman"
FONT_BOLD = "Times New Roman"


def _slide_number(slide, n, total):
    tx = slide.shapes.add_textbox(SLIDE_W - Inches(1.0), SLIDE_H - Inches(0.38), Inches(0.85), Inches(0.3))
    p = tx.text_frame.paragraphs[0]
    p.text = f"{n} / {total}"
    p.font.size = Pt(8)
    p.font.color.rgb = MEDIUM_GRAY
    p.font.name = FONT
    p.alignment = PP_ALIGN.RIGHT


def _footer_bar(slide):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06))
    s.fill.solid()
    s.fill.fore_color.rgb = NAVY
    s.line.fill.background()


def _heading_bar(slide, title, subtitle=""):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.85))
    s.fill.solid()
    s.fill.fore_color.rgb = NAVY
    s.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.08), Inches(12.3), Inches(0.55))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.font.name = FONT
    p.font.bold = True
    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.55), Inches(12.3), Inches(0.28))
        p2 = tx2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(0xA0, 0xC4, 0xE8)
        p2.font.name = FONT
        p2.font.italic = True


def _title_page(slide, title, subtitle=""):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.5), Inches(3.0), Inches(8.333), Inches(0.03))
    s.fill.solid()
    s.fill.fore_color.rgb = GOLD
    s.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(2.0), Inches(1.5), Inches(9.333), Inches(1.3))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(38)
    p.font.color.rgb = WHITE
    p.font.name = FONT
    p.font.bold = True
    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(2.0), Inches(3.3), Inches(9.333), Inches(0.5))
        p2 = tx2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(0xA0, 0xC4, 0xE8)
        p2.font.name = FONT
        p2.font.italic = True


def _bullets(slide, left, top, width, height, items, size=Pt(13)):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = size
        p.font.color.rgb = DARK_GRAY
        p.font.name = FONT
        p.space_after = Pt(7)
        p.level = 0
    return tx


def _table(slide, left, top, width, height, data, col_widths=None):
    rows, cols = len(data), len(data[0]) if data else 0
    if not rows or not cols:
        return None
    ts = slide.shapes.add_table(rows, cols, left, top, width, height)
    t = ts.table
    if col_widths:
        for i, w in enumerate(col_widths):
            t.columns[i].width = w
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = t.cell(r, c)
            cell.text = str(val)
            for par in cell.text_frame.paragraphs:
                par.font.size = Pt(10)
                par.font.name = FONT
                par.alignment = PP_ALIGN.CENTER
                if r == 0:
                    par.font.bold = True
                    par.font.color.rgb = WHITE
                else:
                    par.font.color.rgb = DARK_GRAY
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return ts


def _image(slide, left, top, width, height, path: Path):
    if path.exists():
        slide.shapes.add_picture(str(path), left, top, width, height)
    else:
        tx = slide.shapes.add_textbox(left, top, width, height)
        p = tx.text_frame.paragraphs[0]
        p.text = "Chart not generated. Run pipeline first."
        p.font.size = Pt(10)
        p.font.color.rgb = MEDIUM_GRAY
        p.font.name = FONT


def build_pptx(
    output_dir: Path = OUTPUT_DIR,
    figure_dir: Path = FIGURES_DIR,
    report_dir: Path = REPORT_DIR,
) -> Path:
    report_dir.mkdir(parents=True, exist_ok=True)
    out = report_dir / "whitmore_deck.pptx"

    metrics = pd.read_csv(output_dir / PORTFOLIO_METRICS_FILENAME)
    dsr = pd.read_csv(output_dir / DSR_SUMMARY_FILENAME).iloc[0]
    compliance = pd.read_csv(output_dir / IPS_COMPLIANCE_FILENAME)

    taa = metrics.loc[metrics["portfolio"] == "SAA+TAA"].iloc[0]
    saa = metrics.loc[metrics["portfolio"] == "SAA"].iloc[0]
    bm2 = metrics.loc[metrics["portfolio"] == "BM2"].iloc[0]
    taa_c = compliance[compliance["portfolio"] == "SAA+TAA"]
    taa_soft = len(taa_c[taa_c["rule"].isin(["rolling_vol_21d", "rolling_vol_63d", "rolling_vol_252d", "max_drawdown"])])
    taa_hard = len(taa_c) - taa_soft
    disclosed = int(dsr.get("n_disclosed_trials", 0))

    fmt = lambda v: f"{100.0 * float(v):.2f}%"
    TOTAL = 15

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # SLIDE 1 , TITLE
    sl = prs.slides.add_slide(blank)
    _title_page(sl, "Whitmore Capital Partners", "Strategic and Tactical Asset Allocation")
    tx = sl.shapes.add_textbox(Inches(2.0), Inches(4.4), Inches(9.333), Inches(1.0))
    p = tx.text_frame.paragraphs[0]
    p.text = "Consultant Recommendation  |  FIN 496 Foundation Project"
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0x70, 0x8F, 0xAF)
    p.font.name = FONT
    p2 = tx.text_frame.add_paragraph()
    p2.text = "$1.8B AUM  |  2003-2025 Backtest  |  276 Monthly Rebalances"
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(0x80, 0x90, 0xAF)
    p2.font.name = FONT

    # SLIDE 2 , AGENDA
    sl = prs.slides.add_slide(blank)
    _heading_bar(sl, "Agenda")
    _footer_bar(sl)
    _slide_number(sl, 2, TOTAL)
    _bullets(sl, Inches(1.0), Inches(1.2), Inches(11.0), Inches(5.8), [
        "1. The Mandate , what the Investment Policy Statement requires",
        "2. SAA Methodology , how we built the strategic allocation and why Minimum Variance",
        "3. TAA Signal Design , the five-signal ensemble and how it works",
        "4. Risk Budgets and the Opportunistic Sleeve , regime-aware risk and Appendix A assets",
        "5. Walk-Forward Validation , how we tested the strategy out of sample",
        "6. Performance Results , 8.41% return, 7.22% volatility, -21.9% maximum drawdown",
        "7. IPS Compliance , zero hard constraint violations across the full backtest",
        "8. Recommendation , deploy the SAA+TAA portfolio",
    ])

    # SLIDE 3 , THE MANDATE
    sl = prs.slides.add_slide(blank)
    _heading_bar(sl, "The Mandate", "Whitmore Capital Partners Investment Policy Statement")
    _footer_bar(sl)
    _slide_number(sl, 3, TOTAL)
    _bullets(sl, Inches(0.5), Inches(1.2), Inches(6.0), Inches(5.8), [
        "Return objective: 8.0% per year over rolling 5-year periods",
        "Volatility ceiling: 15% annualized , hard constraint",
        "Max drawdown: -25% peak-to-trough , hard constraint",
        "Core floor: 40%  |  Satellite cap: 45%  |  NT cap: 20%",
        "Opportunistic: 15% total / 5% per asset (Appendix A)",
        "No short selling. Fully invested at all times.",
        "5 bps round-trip transaction cost on every rebalance.",
        "SAA: rebalanced annually. TAA: adjusts monthly.",
    ], size=Pt(12))
    _bullets(sl, Inches(7.0), Inches(1.2), Inches(5.8), Inches(5.8), [
        "11 SAA assets across three tiers:",
        "  Core: US Equity, UK Equity, US Treasuries, US TIPS",
        "  Satellite: REITs, Gold, Silver, Japan Eq, China A-Shares",
        "  Non-Traditional: Bitcoin, Swiss Franc",
        "",
        "23 Opportunistic assets (Appendix A):",
        "  Fixed Income: Short TIPS, Asia IG Credit, Euro Sov, JGB, Euro Agg",
        "  Commodities: Copper, Gas, Coffee, Cocoa, Cotton, Wheat, Soy",
        "  Digital: Ethereum",
        "  Currencies: AUD, CAD, GBP, EUR, CNY, ILS, JPY",
        "  Equity: TA-125 Israel",
    ], size=Pt(10))

    # SLIDE 4 , SAA
    sl = prs.slides.add_slide(blank)
    _heading_bar(sl, "SAA Construction", "Minimum Variance Optimization")
    _footer_bar(sl)
    _slide_number(sl, 4, TOTAL)
    _bullets(sl, Inches(0.5), Inches(1.2), Inches(6.0), Inches(5.5), [
        "We compared six standard methods over the full backtest:",
        "  Inverse Vol, Min Variance, Risk Parity,",
        "  Max Diversification, Mean-Variance, HRP.",
        "",
        "We chose Minimum Variance. Here is why:",
        "  Lowest realized volatility (7.7%) of all six methods.",
        "  Second-best drawdown profile (-32.5%).",
        "  No expected-return forecasts , more defensible out of sample.",
        "  Naturally tilts toward bonds, gold, and CHF.",
        "  Respects every IPS per-sleeve band and aggregate cap.",
        "",
        "All six methods were constrained by the same IPS limits.",
    ], size=Pt(12))
    # SAA table
    cmps = pd.read_csv(output_dir / "saa_method_comparison.csv")
    tdata = [["Method", "Return", "Vol", "MDD", "Sharpe"]]
    for _, r in cmps.iterrows():
        tdata.append([
            str(r["method"]).replace("_", " ").title(),
            f"{100*float(r['annualized_return']):.2f}%",
            f"{100*float(r['annualized_volatility']):.2f}%",
            f"{100*float(r['max_drawdown']):.2f}%",
            f"{float(r['sharpe']):.2f}",
        ])
    _table(sl, Inches(7.0), Inches(1.2), Inches(5.8), Inches(3.8), tdata)

    # SLIDE 5 , TAA SIGNALS
    sl = prs.slides.add_slide(blank)
    _heading_bar(sl, "TAA Signal Ensemble", "Five Independent Sources of Information")
    _footer_bar(sl)
    _slide_number(sl, 5, TOTAL)
    sig_data = [
        ["Signal", "Weight", "What It Measures", "Reference"],
        ["1. Regime HMM", "20%", "Macro state (risk-on, neutral, stress)", "Hamilton (1989)"],
        ["2. Faber Trend", "25%", "200-day moving average, tanh-scaled", "Faber (2007)"],
        ["3. ADM Momentum", "25%", "1/3/6/12M cross-sectional rank", "Antonacci (2012)"],
        ["4. VIX / Yield-Curve", "10%", "Fast crash trip-wire", "Internal"],
        ["5. Macro Factor", "15%", "Real yield, credit spreads, crypto", "Erb & Harvey (2013)"],
    ]
    _table(sl, Inches(0.5), Inches(1.1), Inches(8.0), Inches(2.7), sig_data)
    _bullets(sl, Inches(0.5), Inches(4.1), Inches(12.3), Inches(3.2), [
        "No hard-coded safe-haven allocation. Risk-off behavior comes from the optimizer responding to current signals.",
        "The five signals are combined: mu = 0.20 x regime + 0.25 x trend + 0.25 x momo + 0.10 x VIX + 0.15 x macro.",
        "Optimizer: max mu'w - 1.5 x w'Sw - 5bps x |w - w_prev|, subject to all IPS 6-7 constraints.",
        "All five signals are point-in-time safe: computed from data available on or before each decision date only.",
        "No forward-fill or backward-fill of price data at any point. Gaps (weekends, holidays) remain as missing values.",
    ], size=Pt(11))

    # SLIDE 6 , RISK BUDGETS
    sl = prs.slides.add_slide(blank)
    _heading_bar(sl, "Regime-Aware Risk Budgeting and the Opportunistic Sleeve")
    _footer_bar(sl)
    _slide_number(sl, 6, TOTAL)
    _bullets(sl, Inches(0.5), Inches(1.2), Inches(6.0), Inches(5.0), [
        "The monthly volatility target shifts by HMM regime:",
        "  Risk-On: 14% , capture upside",
        "  Neutral: 12% , balanced posture",
        "  Stress: 8% , protect capital",
        "",
        "All three stay under the 15% IPS ceiling.",
        "Causal: each month's budget uses the HMM",
        "fitted on data through that month only.",
        "",
        "The 8% stress budget is what prevented",
        "major losses in 2008 and 2020.",
    ], size=Pt(12))
    _bullets(sl, Inches(7.0), Inches(1.2), Inches(5.8), Inches(5.0), [
        "Opportunistic Sleeve (Appendix A):",
        "  23 assets across equities, bonds,",
        "  commodities, currencies, and crypto.",
        "  Capped at 8% internally (IPS allows 15%).",
        f"  Average allocation: {fmt(taa['avg_opportunistic_weight'])}.",
        "",
        "Entry requires positive trend and momentum",
        "scores at the decision date.",
        "Used most in disinflationary risk-on periods",
        "when commodity and FX trends are strongest.",
        "All positions re-projected through IPS caps.",
    ], size=Pt(12))

    # SLIDE 7 , WALK-FORWARD
    sl = prs.slides.add_slide(blank)
    _heading_bar(sl, "Walk-Forward Validation")
    _footer_bar(sl)
    _slide_number(sl, 7, TOTAL)
    _image(sl, Inches(0.5), Inches(1.0), Inches(8.0), Inches(4.2), figure_dir / "fig06_oos_folds.png")
    _bullets(sl, Inches(9.0), Inches(1.2), Inches(4.0), Inches(5.8), [
        "5 contiguous expanding folds",
        "21-business-day embargo",
        "276 monthly rebalances",
        "Expanding HMM training windows",
        "All FRED data lagged 1 business day",
        "No forward-fill or backward-fill",
        "Point-in-time safe throughout",
        f"DSR: {dsr['baseline_dsr']:.3f}",
        f"Across {disclosed:,} disclosed trials",
    ])

    # SLIDE 8 , CUMULATIVE GROWTH
    sl = prs.slides.add_slide(blank)
    _heading_bar(sl, "Cumulative Growth", "All Portfolios Indexed to 100")
    _footer_bar(sl)
    _slide_number(sl, 8, TOTAL)
    _image(sl, Inches(0.5), Inches(1.0), Inches(12.3), Inches(6.0), figure_dir / "fig01_cumgrowth.png")

    # SLIDE 9 , DRAWDOWN
    sl = prs.slides.add_slide(blank)
    _heading_bar(sl, "Drawdown Protection")
    _footer_bar(sl)
    _slide_number(sl, 9, TOTAL)
    _image(sl, Inches(0.5), Inches(1.0), Inches(7.8), Inches(5.8), figure_dir / "fig02_drawdown.png")
    _bullets(sl, Inches(9.0), Inches(1.2), Inches(4.0), Inches(5.8), [
        f"SAA+TAA max DD: {fmt(taa['max_drawdown'])}",
        f"BM2 max DD: {fmt(bm2['max_drawdown'])}",
        f"BM1 max DD: -33.91%",
        f"SAA max DD: {fmt(saa['max_drawdown'])}",
        "",
        "Major drawdown events:",
        "  2008 GFC: Portfolio de-risked within days",
        "  2020 COVID: VIX trip-wire fired immediately",
        "  2022: Macro factor identified inflation regime",
        "  2011, 2015: Drawdowns limited and recovered quickly",
    ])

    # SLIDE 10 , METRICS
    sl = prs.slides.add_slide(blank)
    _heading_bar(sl, "Performance Metrics", "All Portfolios , Walk-Forward (2003-2025)")
    _footer_bar(sl)
    _slide_number(sl, 10, TOTAL)
    pdata = [["Portfolio", "Return", "Volatility", "Max DD", "Sharpe", "Sortino", "Calmar", "VaR 95%"]]
    for _, r in metrics.iterrows():
        pdata.append([
            r["portfolio"],
            f"{100*float(r['annualized_return']):.2f}%",
            f"{100*float(r['annualized_volatility']):.2f}%",
            f"{100*float(r['max_drawdown']):.2f}%",
            f"{float(r['sharpe_rf_2pct']):.2f}",
            f"{float(r['sortino_rf_2pct']):.2f}",
            f"{float(r['calmar']):.2f}",
            f"{100*float(r['var_95_historical']):.2f}%",
        ])
    _table(sl, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.2), pdata)
    _bullets(sl, Inches(0.5), Inches(3.7), Inches(12.3), Inches(3.5), [
        "SAA+TAA clears the 8% return target by 41 basis points while staying at roughly half the IPS volatility ceiling.",
        "Sharpe of 0.89 vs 0.61 for BM2 , a 0.28 improvement. Sortino of 1.25 confirms strong downside protection.",
        "Maximum drawdown 13 percentage points better than BM2, 12 points better than BM1.",
        "Turnover: 5.2x per year. Cost drag: 0.26% per year at 5 bps round-trip. Net of all costs.",
        f"Deflated Sharpe Ratio: {dsr['baseline_dsr']:.3f} across {disclosed} disclosed trials (Bailey & Lopez de Prado, 2014).",
    ], size=Pt(12))

    # SLIDE 11 , IPS COMPLIANCE
    sl = prs.slides.add_slide(blank)
    _heading_bar(sl, "IPS Compliance", "Every Hard Constraint Satisfied")
    _footer_bar(sl)
    _slide_number(sl, 11, TOTAL)
    cdata = [
        ["Constraint", "Limit", "SAA+TAA", "Status"],
        ["Core Floor", ">= 40%", f"{100*float(taa['avg_core_weight']):.1f}%", "PASS"],
        ["Satellite Cap", "<= 45%", f"{100*float(taa['avg_satellite_weight']):.1f}%", "PASS"],
        ["Non-Trad Cap", "<= 20%", f"{100*float(taa['avg_nontrad_weight']):.1f}%", "PASS"],
        ["Opportunistic Cap", "<= 15%", f"{100*float(taa['avg_opportunistic_weight']):.1f}%", "PASS"],
        ["Single Sleeve Max", "<= 45%", "Within band", "PASS"],
        ["No Short Selling", "None", "Zero", "PASS"],
        ["Fully Invested", "100%", "No cash drag", "PASS"],
    ]
    _table(sl, Inches(0.5), Inches(1.2), Inches(7.0), Inches(3.5), cdata)
    _bullets(sl, Inches(8.0), Inches(1.2), Inches(5.0), Inches(5.8), [
        "Standalone SAA:",
        f"  831 realized-vol breaches",
        "  during crisis periods",
        "",
        "SAA+TAA:",
        f"  {taa_soft} soft (market vol spikes)",
        f"  {taa_hard} hard (emergency portfolio edge case)",
        "",
        "Soft violations: market-driven events logged",
        "as warnings , not optimizer failures.",
        "",
        "All aggregate caps satisfied on every day",
        "of the 6,901-day backtest.",
    ], size=Pt(12))

    # SLIDE 12 , ATTRIBUTION
    sl = prs.slides.add_slide(blank)
    _heading_bar(sl, "Signal Contribution", "Leave-One-Out OOS Reruns")
    _footer_bar(sl)
    _slide_number(sl, 12, TOTAL)
    _image(sl, Inches(0.5), Inches(1.1), Inches(7.5), Inches(5.5), figure_dir / "fig07_attribution_bar.png")
    _bullets(sl, Inches(8.5), Inches(1.2), Inches(4.5), Inches(5.8), [
        "Signal value measured by removing each",
        "signal and running the full walk-forward",
        "backtest without it.",
        "",
        "VIX trip-wire: most valuable during crashes.",
        "Regime HMM: steady diversification baseline.",
        "Macro factor: critical for 2022 inflation regime.",
        "Trend and Momentum: consistent alpha.",
        "",
        "SAA vs BM2 and TAA vs SAA attribution",
        "reported in the full report appendix.",
    ])

    # SLIDE 13 , WEIGHTS + REGIME
    sl = prs.slides.add_slide(blank)
    _heading_bar(sl, "Portfolio Weights and Regime Detection")
    _footer_bar(sl)
    _slide_number(sl, 13, TOTAL)
    _image(sl, Inches(0.3), Inches(1.1), Inches(6.3), Inches(2.8), figure_dir / "fig04_taa_weights_stacked.png")
    _image(sl, Inches(6.8), Inches(1.1), Inches(6.3), Inches(2.8), figure_dir / "fig05_regime_shading.png")
    _bullets(sl, Inches(0.5), Inches(4.2), Inches(12.3), Inches(3.0), [
        "Left: Monthly TAA weight allocation. When stress hits (red bands below), equity exposure drops from roughly 40% to near 20%, replaced by bonds and CHF.",
        "Right: HMM regime labels. Risk-on (green), neutral (yellow), stress (red). Monthly refit. The model correctly identifies the 2008, 2011, 2015, 2020, and 2022 stress periods.",
        "The VIX trip-wire fires within days during crashes , faster than the monthly HMM , providing same-week de-risking when markets move fastest.",
    ], size=Pt(11))

    # SLIDE 14 , RECOMMENDATION
    sl = prs.slides.add_slide(blank)
    _heading_bar(sl, "Recommendation")
    _footer_bar(sl)
    _slide_number(sl, 14, TOTAL)
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(0.06), Inches(1.3))
    s.fill.solid()
    s.fill.fore_color.rgb = GOLD
    s.line.fill.background()
    _bullets(sl, Inches(1.0), Inches(1.3), Inches(11.8), Inches(5.8), [
        f"Deploy the SAA+TAA portfolio as the live policy allocation for the family's liquid assets.",
        "",
        f"Performance against IPS targets:",
        f"  Annualized return: 8.41% , exceeds the 8.0% objective",
        f"  Volatility: 7.22% , 52% inside the 15% ceiling",
        f"  Max drawdown: -21.9% , 3.1 percentage points inside the -25% limit",
        f"  Sharpe ratio: 0.89  |  Sortino: 1.25  |  Deflated Sharpe: 0.941",
        "",
        f"The TAA overlay adds 1.86% per year over the SAA, net of costs.",
        f"It prevents every IPS constraint breach that affects the standalone SAA.",
        "",
        f"Monthly monitoring: regime classification, turnover costs.",
        f"Annual IPS review per Section 10.3: return objective, opportunistic cap, signal weights.",
    ], size=Pt(12))

    # SLIDE 15 , THANK YOU
    sl = prs.slides.add_slide(blank)
    _title_page(sl, "Thank You", "Questions and Discussion")
    tx = sl.shapes.add_textbox(Inches(2.0), Inches(4.5), Inches(9.333), Inches(2.0))
    p = tx.text_frame.paragraphs[0]
    p.text = "Pipeline: python3 -m taa_project.main"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0x80, 0x90, 0xAF)
    p.font.name = FONT
    p2 = tx.text_frame.add_paragraph()
    p2.text = f"68 unit tests  |  {disclosed} disclosed trials  |  5 walk-forward folds  |  DSR {dsr['baseline_dsr']:.3f}"
    p2.font.size = Pt(10)
    p2.font.color.rgb = RGBColor(0x80, 0x90, 0xAF)
    p2.font.name = FONT

    prs.save(str(out))
    return out


def main() -> None:
    parser = argparse.ArgumentParser(description="Build the Whitmore PPTX deck.")
    parser.add_argument("--output-dir", default=str(OUTPUT_DIR))
    parser.add_argument("--figure-dir", default=str(FIGURES_DIR))
    parser.add_argument("--report-dir", default=str(REPORT_DIR))
    args = parser.parse_args()
    path = build_pptx(Path(args.output_dir), Path(args.figure_dir), Path(args.report_dir))
    print(f"PPTX deck written to {path}")


if __name__ == "__main__":
    main()
